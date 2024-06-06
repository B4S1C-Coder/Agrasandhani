import json
from PyPDF2 import PdfReader
from pptx import Presentation

from langchain.llms import Anthropic

class Claude2AnthropicLargeLanguageModel:
    def __init__(self, max_tokens=700, *args, **kwargs):

        self.__llm_instance = Anthropic(model="claude-2",
                            max_tokens=max_tokens, *args, **kwargs)

        self.question_generation_prompt = """
Generate 5 multiple choice questions from the following text.
All the questions must be enclosed in a single JSON object. 
Further each question will be a JSON object in the format:
question: question, 1: option1, 2: option2, 3:option3, 4:option4, ans: correctoption\n
Text: """

    def __determinePDForPPT(self, file_obj):
        # PDF file has first 5 bytes as %PDF-
        file_signature_PDF = file_obj.read(5)
        # Reset the cursor
        file_obj.seek(0)
        # PPT file is a zip file and zip file has first two bytes as PK
        file_signature_PPT = file_obj.read(2)

        if file_signature_PDF == b'%PDF-':
            return "PDF"
        
        if file_signature_PPT == b'PK':
            return "PPTX"

        return None

    def generate_questions_from_text(self, text):
        prompt = self.question_generation_prompt + text

        response = self.__llm_instance.generate(prompts=[prompt])
        response_text = response.generations[0][0].text

        try:
            response_json = json.loads(response_text)
            return dict(response_json)
        except Exception as err:
            print("Failed to convert to JSON.")
            print(f"[ ERROR ] {err}")
            return response_text

    def extract_text_from_file(self, file_obj):
        file_type = self.__determinePDForPPT(file_obj)
        file_obj.seek(0) # Just in case

        if file_type == "PDF":
            reader = PdfReader(file_obj)

            # print(f"[ LLM - DEBUG ] {len(file_obj)}", file="logs.txt")

            extracted_text = []

            def ignore_header_footer(text, cm, tm, fontDict, fontSize):
                # global extracted_text
                y = tm[5]

                if y > 50 and y < 720:
                    # extracted_text += text
                    extracted_text.append(text)

            for page in reader.pages:
                page.extract_text(visitor_text=ignore_header_footer)

            extracted_text = "".join(extracted_text)

            return extracted_text

        if file_type == "PPTX":
            my_ppt = Presentation(file_obj)

            extracted_text = ""

            for slide in my_ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text

            return extracted_text

        return ""